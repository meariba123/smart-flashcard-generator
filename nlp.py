import re
import random
import os
import spacy
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
import pytesseract
from rapidfuzz import fuzz
from PIL import Image, ImageDraw, ImageFont
import textwrap

# LOAD SPACY
try:
    nlp = spacy.load("en_core_web_sm")
except:
    os.system("python -m spacy download en_core_web_sm")
    nlp = spacy.load("en_core_web_sm")

# CONSTANTS
EXACT_JUNK = {
    'this', 'that', 'it', 'they', 'there',
    'what', 'which', 'who', 'example', 'examples'
}

URL_PATTERN = r'(https?://\S+|www\.\S+|\S+\.com/\S+|\S+\.be/\S+)'
PREFIX_JUNK = r'^(e\.g\.|i\.e\.|etc|example:|note:)\s+'
MAX_VISUALS = 3

# CLEANING FUNCTIONS
def clean_term(term):
    term = re.sub(PREFIX_JUNK, '', term, flags=re.IGNORECASE)
    term = re.sub(URL_PATTERN, '', term)
    return term.strip().capitalize()

def clean_text(text):
    text = re.sub(URL_PATTERN, '', text)
    text = re.sub(r'\(.*?\)', '', text)
    text = text.replace(' - ', ' ').replace(' : ', ' ')
    return text.strip()

def is_valid_term(term):
    if not term: return False
    term_lower = term.lower()
    if term_lower in EXACT_JUNK: return False
    if any(x in term_lower for x in ["http", "www", "//", ".com", ".be"]): return False
    if len(term.split()) > 7: return False
    if len(term) < 3: return False
    if not re.search(r'[a-zA-Z]', term): return False
    return True

# STRUCTURED VISUAL EXPLANATION
def generate_visual_explanation(term):
    try:
        safe_term = re.sub(r'[^a-zA-Z0-9_]', '_', term.strip())
        img = Image.new("RGB", (900, 450), "white")
        draw = ImageDraw.Draw(img)

        try:
            font_title = ImageFont.truetype("arial.ttf", 42)
            font_body = ImageFont.truetype("arial.ttf", 22)
        except:
            font_title = ImageFont.load_default()
            font_body = ImageFont.load_default()

        doc = nlp(term)
        keywords = [token.text for token in doc if token.pos_ in ["NOUN", "PROPN"]]

        concept = term
        process = f"Involves {', '.join(keywords[:3])}" if keywords else "Core mechanism"
        outcome = "Produces structured output"

        draw.text((40, 20), term, fill="black", font=font_title)
        
        concept_wrapped = textwrap.fill(concept, width=28)
        process_wrapped = textwrap.fill(process, width=28)
        outcome_wrapped = textwrap.fill(outcome, width=28)

        draw.text((60, 150), f"Concept:\n{concept_wrapped}", fill="blue", font=font_body)
        draw.text((350, 150), f"Process:\n{process_wrapped}", fill="green", font=font_body)
        draw.text((650, 150), f"Outcome:\n{outcome_wrapped}", fill="red", font=font_body)

        os.makedirs("static/generated_images", exist_ok=True)
        image_path = f"static/generated_images/{safe_term}.png"
        img.save(image_path)
        return f"generated_images/{safe_term}.png"
    except Exception as e:
        print("VISUAL EXPLANATION FAILED:", str(e))
        return None

# TEXT EXTRACTION
def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    text = ""
    try:
        if ext == ".txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == ".docx":
            doc = Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".pdf":
            reader = PdfReader(filepath)
            text = "\n".join([page.extract_text() or "" for page in reader.pages])
        elif ext in [".png", ".jpg", ".jpeg"]:
            text = pytesseract.image_to_string(Image.open(filepath))
        elif ext == ".pptx":
            prs = Presentation(filepath)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        full_text.append(shape.text)
            text = "\n".join(full_text)
    except Exception as e:
        print(f"Error extracting text: {e}")
    return text

# ----------------------------
# FLASHCARD GENERATION (CORE UPDATE)
# ----------------------------
def generate_flashcards_from_file(filepath, language='en'):
    ext = os.path.splitext(filepath)[1].lower()
    flashcards = []

    # Language Config
    lang_config = {
        'es': {"prefix": "¿Qué es", "suffix": "?"},
        'en': {"prefix": "What is", "suffix": "?"}
    }
    l = lang_config.get(language, lang_config['en'])

    if ext == ".pptx":
        prs = Presentation(filepath)
        for slide in prs.slides:
            title, content_lines = "", []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if not title: title = text
                    else: content_lines.append(text)

            if not title or len(title) > 80: continue
            combined_content = " ".join(content_lines)
            if len(combined_content) < 20: continue

            if is_valid_term(title):
                flashcards.append({
                    "question": f"{l['prefix']} {title.strip()}{l['suffix']}",
                    "answer": combined_content.strip(),
                    "visual_explanation": None,
                    "score": 0.9,
                    "language": language
                })
    else:
        text = extract_text_from_file(filepath)
        lines = text.split('\n')
        for line in lines:
            line = line.strip()
            if ":" in line:
                parts = line.split(":", 1)
                term = clean_term(parts[0])
                definition = clean_text(parts[1])
                if is_valid_term(term) and len(definition) > 15:
                    flashcards.append({
                        "question": f"{l['prefix']} {term}{l['suffix']}",
                        "answer": definition,
                        "visual_explanation": None,
                        "score": 0.8,
                        "language": language
                    })

    # DEDUPLICATE
    unique = {}
    for card in flashcards:
        unique[card["question"].lower()] = card
    flashcards = list(unique.values())

    if not flashcards: return []

    # GENERATE VISUALS (stripping lang prefixes for image generator)
    for card in flashcards[:MAX_VISUALS]:
        term_to_draw = card["question"]
        for p in ["What is ", "¿Qué es "]: term_to_draw = term_to_draw.replace(p, "")
        term_to_draw = term_to_draw.replace("?", "")
        card["visual_explanation"] = generate_visual_explanation(term_to_draw)

    random.shuffle(flashcards)
    return flashcards

# ANSWER CHECKING
def is_answer_correct(user_answer, correct_answer, threshold=75):
    user_answer = user_answer.lower().strip()
    correct_answer = correct_answer.lower().strip()
    similarity = fuzz.ratio(user_answer, correct_answer)
    if similarity >= threshold: return True
    user_keywords = set(user_answer.split())
    correct_keywords = set(correct_answer.split())
    if len(correct_keywords) > 0:
        overlap = len(user_keywords & correct_keywords) / len(correct_keywords)
        if overlap > 0.6: return True
    return False